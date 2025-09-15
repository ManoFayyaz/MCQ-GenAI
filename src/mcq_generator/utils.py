#helper file: helpeing functions
import os
import json
import re
import streamlit as st
from PyPDF2 import PdfReader
import traceback
from pptx import Presentation
from PIL import Image
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"D:\Tesseract-OCR\tesseract.exe"
from pdf2image import convert_from_bytes
from src.mcq_generator.logger import logging as mcq_logger
import io
from sentence_transformers import SentenceTransformer
import chromadb

# Load embedding model once
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

# Initialize Chroma client (local persistence)
chroma_client = chromadb.PersistentClient(path="chroma_store")
collection = chroma_client.get_or_create_collection(name="documents")

# --------- Helper: split text into chunks ----------
def split_text(text, chunk_size=500, overlap=50):
    """
    Splits text into smaller chunks for embedding.
    Example: chunk_size=500 chars, with 50 overlap.
    """
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - overlap  # step back a little to maintain context
    return chunks


def read_file(file, file_id="default_doc"):
    text = ""

    # ---- PDF ----
    if file.name.endswith('.pdf'):
        try:
            pdf_reader = PdfReader(file)
            images = convert_from_bytes(
                file.getbuffer(),
                poppler_path=r"D:\Poppler\poppler-25.07.0\Library\bin"
            )

            for page_num, page in enumerate(pdf_reader.pages):
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"

                if page_num < len(images):
                    img = images[page_num]
                    text += pytesseract.image_to_string(img) + "\n"

        except Exception as e:
            mcq_logger.error("Error reading PDF file: %s", e)
            st.error(f"An error occurred while reading PDF: {e}")

    # ---- TXT ----
    elif file.name.endswith('.txt'):
        try:
            text = file.read().decode('utf-8')
        except Exception as e:
            mcq_logger.error("Error reading TXT file: %s", e)
            st.error(f"An error occurred while reading TXT: {e}")

    # ---- PPT / PPTX ----
    elif file.name.endswith(('.ppt', '.pptx')):
        try:
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                    elif shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text += cell.text + "\n"
                    elif shape.shape_type == 13:  # Picture
                        image = shape.image
                        image_bytes = image.blob
                        img = Image.open(io.BytesIO(image_bytes))
                        text += pytesseract.image_to_string(img) + "\n"

        except Exception as e:
            mcq_logger.error("Error reading PPT file: %s", e)
            st.error(f"An error occurred while reading PPT: {e}")

    # Store embeddings in Chroma
    
    # ---- Store in Chroma ----
    if text.strip():
        chunks = split_text(text, chunk_size=500, overlap=50)
        embeddings = embedding_model.encode(chunks).tolist()

        # Store each chunk separately in Chroma
        ids = [f"{file_id}_{i}" for i in range(len(chunks))]
        collection.upsert(
            documents=chunks,
            embeddings=embeddings,
            ids=ids
        )

    return text.strip()

# print(chroma_client.list_collections())

# def read_file(file):
#     text = ""

#     # ---- PDF ----
#     if file.name.endswith('.pdf'):
#         try:
            
#             pdf_reader = PdfReader(file)
#             images =   images = convert_from_bytes(
#                 file.getbuffer(),
#                 poppler_path=r"D:\Poppler\poppler-25.07.0\Library\bin"  # <-- specify Poppler path here
#             ) 

#             for page_num, page in enumerate(pdf_reader.pages):
#                 # 1. Extract normal text
#                 extracted = page.extract_text()
#                 if extracted:
#                     text += extracted + "\n"

#                 # 2. OCR from images (always run OCR too)
#                 if page_num < len(images):
#                     img = images[page_num]
#                     text += pytesseract.image_to_string(img) + "\n"

#             return text.strip()

#         except Exception as e:
#             mcq_logger.error("Error reading PDF file: %s", e)
#             st.error(f"An error occurred while reading PDF: {e}")
            
#     # ---- TXT ----
#     elif file.name.endswith('.txt'):
#         try:
#             text = file.read().decode('utf-8')
#             return text.strip()
#         except Exception as e:
#             mcq_logger.error("Error reading TXT file: %s", e)
#             st.error(f"An error occurred while reading TXT: {e}")

#     # ---- PPT / PPTX ----
#     elif file.name.endswith('.ppt') or file.name.endswith('.pptx'):
#         try:
#             prs = Presentation(file)

#             for slide in prs.slides:
#                 for shape in slide.shapes:
#                     # 1. Extract text
#                     if hasattr(shape, "text") and shape.text:
#                         text += shape.text + "\n"

#                     # 2. Extract table text
#                     elif shape.has_table:
#                         for row in shape.table.rows:
#                             for cell in row.cells:
#                                 text += cell.text + "\n"

#                     # 3. OCR from images
#                     elif shape.shape_type == 13:  # Picture
#                         image = shape.image
#                         image_bytes = image.blob
#                         img = Image.open(io.BytesIO(image_bytes))
#                         text += pytesseract.image_to_string(img) + "\n"

#             return text.strip()

#         except Exception as e:
#             mcq_logger.error("Error reading PPT file: %s", e)
#             st.error(f"An error occurred while reading PPT: {e}")

#     return text.strip()



def get_table_data(quiz):
    try:
        quiz_clean = re.sub(r"^```json|```$", "", quiz, flags=re.MULTILINE).strip()

        quiz_dict = json.loads(quiz_clean)
        
        quiz_table=[]
        for key, value in quiz_dict.items():
            question = value.get("question", "")
            options = value.get("options", {})
            correct_answer = value.get("correct_answer", "")
            
            quiz_table.append({
                "Question": question,
                "Option A": options.get("A", ""),
                "Option B": options.get("B", ""),
                "Option C": options.get("C", ""),
                "Option D": options.get("D", ""),
                "Correct Answer": correct_answer
            })
        return quiz_table
    
    except json.JSONDecodeError as e:
        print("Error decoding JSON:", e)
        print("Quiz content:", quiz)
        return []
