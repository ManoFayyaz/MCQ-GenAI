# #helper file: helpeing functions
# backend/utils.py
import os
import json
import re
import zipfile
from io import BytesIO
from PyPDF2 import PdfReader
from openai import embeddings
from pptx import Presentation
from PIL import Image
import pytesseract
from pdf2image import convert_from_bytes
from docx import Document
import io
import hashlib
# from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.utils import embedding_functions
from src.mcq_generator.logger import logging as mcq_logger

# Set these to your local locations if needed
pytesseract.pytesseract.tesseract_cmd = r"D:\Tesseract-OCR\tesseract.exe"
POPLER_PATH = r"D:\Poppler\poppler-25.07.0\Library\bin"

# Embedding model and chroma client
# embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
chroma_client = chromadb.EphemeralClient()
embedding_fn = embedding_functions.DefaultEmbeddingFunction()

# chroma_client = chromadb.PersistentClient(path="chroma_store")


def split_text(text, chunk_size=500, overlap=50):
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - overlap
    return chunks

def file_hash(file_obj):
    file_obj.seek(0)
    digest = hashlib.md5(file_obj.read()).hexdigest()
    file_obj.seek(0)
    return digest

def read_file_and_index(file_obj, user_id, file_hash_id):
    """
    Read file bytes, extract text (OCR where needed), embed and upsert into a
    per-(user,file) Chroma collection. Returns (text, collection_name).
    """
    text = ""
    name = getattr(file_obj, "name", "") or "uploaded_file"

    # PDF
    if name.lower().endswith('.pdf'):
        try:
            pdf_reader = PdfReader(file_obj)
            try:
                file_obj.seek(0)
                images = convert_from_bytes(file_obj.getbuffer(), poppler_path=POPLER_PATH)
            except Exception:
                images = []
            for page_num, page in enumerate(pdf_reader.pages):
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
                if page_num < len(images):
                    img = images[page_num]
                    text += pytesseract.image_to_string(img) + "\n"
        except Exception as e:
            mcq_logger.error("Error reading PDF file: %s", e)

    # DOCX
    elif name.lower().endswith('.docx'):
        try:
            doc = Document(file_obj)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
            file_obj.seek(0)
            file_bytes = BytesIO(file_obj.read())
            with zipfile.ZipFile(file_bytes) as docx_zip:
                for file_name in docx_zip.namelist():
                    if file_name.startswith("word/media/"):
                        with docx_zip.open(file_name) as image_file:
                            img = Image.open(image_file)
                            text += pytesseract.image_to_string(img) + "\n"
        except Exception as e:
            mcq_logger.error("Error reading DOCX file: %s", e)

    # TXT
    elif name.lower().endswith('.txt'):
        try:
            file_obj.seek(0)
            raw = file_obj.read()
            text = raw.decode('utf-8', errors='ignore') if isinstance(raw, (bytes, bytearray)) else str(raw)
        except Exception as e:
            mcq_logger.error("Error reading TXT file: %s", e)

    # PPT/PPTX
    elif name.lower().endswith(('.ppt', '.pptx')):
        try:
            prs = Presentation(file_obj)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                    elif getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text += cell.text + "\n"
                    else:
                        try:
                            image = getattr(shape, "image", None)
                            if image is not None:
                                image_bytes = image.blob
                                img = Image.open(io.BytesIO(image_bytes))
                                text += pytesseract.image_to_string(img) + "\n"
                        except Exception:
                            pass
        except Exception as e:
            mcq_logger.error("Error reading PPT file: %s", e)

    # Create or reuse a collection specific for this user-file
    collection_name = f"doc_user{user_id}_{file_hash_id}"
    collection = chroma_client.get_or_create_collection(name=collection_name)

    if text.strip():
        try:
            chunks = split_text(text, chunk_size=500, overlap=50)
            embeddings = embedding_fn(chunks)
            ids = [f"{file_hash_id}_{i}" for i in range(len(chunks))]
            collection.upsert(documents=chunks, embeddings=embeddings, ids=ids)
            # embeddings = embedding_model.encode(chunks).tolist()
            # ids = [f"{file_hash_id}_{i}" for i in range(len(chunks))]
            # collection.upsert(documents=chunks, embeddings=embeddings, ids=ids)
        except Exception as e:
            mcq_logger.error("Error embedding/upserting: %s", e)

    return text.strip(), collection_name

def retrieve_context(query, collection_name, top_k=5):
    try:
        collection = chroma_client.get_collection(name=collection_name)
        results = collection.query(query_texts=[query], n_results=top_k)
        retrieved_chunks = results["documents"][0]
        return " ".join(retrieved_chunks)
    except Exception as e:
        mcq_logger.error("Error retrieving context: %s", e)
        return ""

def normalize_quiz_json(quiz_data):
    """
    Acts as a secure validation safety net for incoming data arrays.
    """
    # 1. If LangChain already delivered a parsed dictionary, return it immediately
    if isinstance(quiz_data, dict):
        return quiz_data
        
    # 2. Safety handler if a raw text string somehow slips through
    if isinstance(quiz_data, str):
        try:
            txt = re.sub(r"^```json|```$", "", quiz_data, flags=re.MULTILINE).strip()
            
            # Isolate json bounds
            start_idx = txt.find('{')
            end_idx = txt.rfind('}')
            if start_idx != -1 and end_idx != -1:
                txt = txt[start_idx:end_idx + 1]
                
            return json.loads(txt)
        except Exception as e:
            mcq_logger.error("Failed to parse fallback quiz string text: %s", e)
            return {}
            
    # 3. If data format is null or unexpected
    return {}

# def normalize_quiz_json(quiz_str):
#     try:
#         txt = re.sub(r"^```json|```$", "", quiz_str, flags=re.MULTILINE).strip()
#         txt = txt.replace('" correct answer"', '"correct_answer"').replace("' correct answer'", "'correct_answer'")
#         obj = json.loads(txt)
#         return obj
#     except Exception:
#         try:
#             txt2 = txt.replace("'", '"')
#             obj = json.loads(txt2)
#             return obj
#         except Exception as e:
#             mcq_logger.error("Failed to parse quiz JSON: %s", e)
#             return {}























# import os
# import json
# import re
# import zipfile
# from io import BytesIO
# import streamlit as st
# from PyPDF2 import PdfReader
# import traceback
# from pptx import Presentation
# from PIL import Image
# import pytesseract
# pytesseract.pytesseract.tesseract_cmd = r"D:\Tesseract-OCR\tesseract.exe"
# from pdf2image import convert_from_bytes
# from docx import Document 
# from src.mcq_generator.logger import logging as mcq_logger
# import io
# from sentence_transformers import SentenceTransformer
# import chromadb
 
# # Load embedding model once
# embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

# # Initialize Chroma client (local persistence)
# chroma_client = chromadb.PersistentClient(path="chroma_store")

# collection = chroma_client.get_or_create_collection(name="documents")

# # --------- Helper: split text into chunks ----------
# def split_text(text, chunk_size=500, overlap=50):
#     """
#     Splits text into smaller chunks for embedding.
#     Example: chunk_size=500 chars, with 50 overlap.
#     """
#     chunks = []
#     start = 0
#     while start < len(text):
#         end = start + chunk_size
#         chunk = text[start:end]
#         chunks.append(chunk)
#         start = end - overlap  # step back a little to maintain context
#     return chunks


# def read_file(file, file_id="default_doc"):
#     text = ""

#     # ---- PDF ----
#     if file.name.endswith('.pdf'):
#         try:
#             pdf_reader = PdfReader(file)
#             images = convert_from_bytes(
#                 file.getbuffer(),
#                 poppler_path=r"D:\Poppler\poppler-25.07.0\Library\bin"
#             )

#             for page_num, page in enumerate(pdf_reader.pages):
#                 extracted = page.extract_text()
#                 if extracted:
#                     text += extracted + "\n"

#                 if page_num < len(images):
#                     img = images[page_num]
#                     text += pytesseract.image_to_string(img) + "\n"

#         except Exception as e:
#             mcq_logger.error("Error reading PDF file: %s", e)
#             st.error(f"An error occurred while reading PDF: {e}")

#     #--- DOCX---
#     elif file.name.endswith('.docx'):
#         try:
#             doc = Document(file)
#             for para in doc.paragraphs:
#                 if para.text.strip():
#                     text += para.text + "\n"

#             file.seek(0)
#             images_text = ""
#             file_bytes = BytesIO(file.read())
#             with zipfile.ZipFile(file_bytes) as docx_zip:
#                 for file_name in docx_zip.namelist():
#                     if file_name.startswith("word/media/"):
#                         with docx_zip.open(file_name) as image_file:
#                             img = Image.open(image_file)
#                             images_text += pytesseract.image_to_string(img) + "\n"
                            
#         except Exception as e:
#             mcq_logger.error("Error reading DOCX file: %s", e)
#             st.error(f"An error occurred while reading DOCX: {e}")

#     # ---- TXT ----
#     elif file.name.endswith('.txt'):
#         try:
#             text = file.read().decode('utf-8')
#         except Exception as e:
#             mcq_logger.error("Error reading TXT file: %s", e)
#             st.error(f"An error occurred while reading TXT: {e}")

#     # ---- PPT / PPTX ----
#     elif file.name.endswith(('.ppt', '.pptx')):
#         try:
#             prs = Presentation(file)
#             for slide in prs.slides:
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text") and shape.text:
#                         text += shape.text + "\n"
#                     elif shape.has_table:
#                         for row in shape.table.rows:
#                             for cell in row.cells:
#                                 text += cell.text + "\n"
#                     elif shape.shape_type == 13:  # Picture
#                         image = shape.image
#                         image_bytes = image.blob
#                         img = Image.open(io.BytesIO(image_bytes))
#                         text += pytesseract.image_to_string(img) + "\n"

#         except Exception as e:
#             mcq_logger.error("Error reading PPT file: %s", e)
#             st.error(f"An error occurred while reading PPT: {e}")

#     # Store embeddings in Chroma
#     # collection_name = f"doc_{user_id}_{file_id}"

#     # # ---- Create or reuse that collection ----
#     # collection = chroma_client.get_or_create_collection(name=collection_name)

#     # ---- Store in Chroma ----
#     if text.strip():
#         try:
#             chroma_client.delete_collection("documents")
#             mcq_logger.info("Old Chroma collection deleted")
#         except Exception as e:
#              mcq_logger.warning(f"Collection didn't exist or couldn't be deleted: {e}")


#         collection = chroma_client.get_or_create_collection(name="documents")

#         chunks = split_text(text, chunk_size=500, overlap=50)
#         embeddings = embedding_model.encode(chunks).tolist()

#         # Store each chunk separately in Chroma
#         ids = [f"{file_id}_{i}" for i in range(len(chunks))]
#         collection.upsert(
#             documents=chunks,
#             embeddings=embeddings,
#             ids=ids
#         )

#     return text.strip(),collection

# def retrieve_context(query,collection,top_k=5):
#     results = collection.query(
#         query_texts=[query],
#         n_results=top_k
#     )
#     # documents come back as list of lists
#     retrieved_chunks = results["documents"][0]
#     return " ".join(retrieved_chunks)


# def get_table_data(quiz):
#     try:
#         quiz_clean = re.sub(r"^```json|```$", "", quiz, flags=re.MULTILINE).strip()

#         quiz_dict = json.loads(quiz_clean)
        
#         quiz_table=[]
#         for key, value in quiz_dict.items():
#             question = value.get("question", "")
#             options = value.get("options", {})
#             correct_answer = value.get("correct_answer", "")
            
#             quiz_table.append({
#                 "Question": question,
#                 "Option A": options.get("A", ""),
#                 "Option B": options.get("B", ""),
#                 "Option C": options.get("C", ""),
#                 "Option D": options.get("D", ""),
#                 "Correct Answer": correct_answer
#             })
#         return quiz_table
    
#     except json.JSONDecodeError as e:
#         print("Error decoding JSON:", e)
#         print("Quiz content:", quiz)
#         return []
